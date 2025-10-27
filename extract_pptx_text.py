from pptx import Presentation
import os

def extract_text_from_pptx(pptx_path, output_file=None):
    """Extract text from a PowerPoint file and optionally save to a text file."""
    prs = Presentation(pptx_path)
    text = []
    
    for i, slide in enumerate(prs.slides, 1):
        text.append(f"\n=== Slayd {i} ===\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    full_text = "\n".join(text)
    
    if output_file:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(full_text)
        return f"Matn faylga saqlandi: {output_file}"
    
    return full_text

if __name__ == "__main__":
    pptx_file = "taksi_loixasi.pptm"
    output_file = "taksi_loixasi_matni.txt"
    
    if os.path.exists(pptx_file):
        result = extract_text_from_pptx(pptx_file, output_file)
        print(result)
    else:
        print(f"Xatolik: {pptx_file} fayli topilmadi!")
